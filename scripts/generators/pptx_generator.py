"""Generate synthetic PowerPoint (.pptx) test files using python-pptx."""

from pathlib import Path

from generators.base import BaseGenerator

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
    _PPTX_OK = True
except ImportError:
    _PPTX_OK = False


def _require_pptx() -> None:
    if not _PPTX_OK:
        raise ImportError("python-pptx is required: uv add python-pptx")


class PptxGenerator(BaseGenerator):
    category = "pptx"

    def generate(self, dry_run: bool = False) -> list[dict]:
        _require_pptx()
        out = self.output_dir()
        if not dry_run:
            out.mkdir(parents=True, exist_ok=True)

        entries = []
        entries.append(self._basic_bullets(out, dry_run))
        entries.append(self._table_slides(out, dry_run))
        entries.append(self._complex_layout(out, dry_run))
        entries.append(self._speaker_notes(out, dry_run))
        return entries

    # ------------------------------------------------------------------
    def _basic_bullets(self, out: Path, dry_run: bool) -> dict:
        filename = "syn_basic_bullets.pptx"
        if dry_run:
            print(f"  [dry-run] would create {self.category}/{filename}")
        else:
            prs = Presentation()
            blank_layout = prs.slide_layouts[6]  # Blank

            # Title slide
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = "Q4 Business Review"
            slide.placeholders[1].text = "Presented by the Finance & Strategy Team\nDecember 2024"

            # Agenda slide
            slide2 = prs.slides.add_slide(prs.slide_layouts[1])
            slide2.shapes.title.text = "Agenda"
            tf = slide2.placeholders[1].text_frame
            tf.text = "1. Financial Performance"
            for item in [
                "2. Market Share Update",
                "3. Key Wins & Challenges",
                "4. Q1 2025 Outlook",
                "5. Q&A",
            ]:
                p = tf.add_paragraph()
                p.text = item
                p.level = 0

            # Slide 3: Financial highlights
            slide3 = prs.slides.add_slide(prs.slide_layouts[1])
            slide3.shapes.title.text = "Financial Performance — Q4 2024"
            tf3 = slide3.placeholders[1].text_frame
            tf3.text = "Revenue: $284M (+14% YoY)"
            for item in [
                "Gross Margin: 41.2% (+2.1pp)",
                "Operating Income: $52M (+22% YoY)",
                "Free Cash Flow: $38M (record high)",
                "Headcount: 1,842 (+6% YoY)",
            ]:
                p = tf3.add_paragraph()
                p.text = item

            prs.save(out / filename)
        return self._entry(filename, "easy",
                           ["pptx", "title-slide", "bullet-points", "financial-review"],
                           False, "Three-slide PPTX: title, agenda, and financial highlights bullets.")

    def _table_slides(self, out: Path, dry_run: bool) -> dict:
        filename = "syn_table_slides.pptx"
        if dry_run:
            print(f"  [dry-run] would create {self.category}/{filename}")
        else:
            prs = Presentation()

            # Slide 1: Title
            slide1 = prs.slides.add_slide(prs.slide_layouts[0])
            slide1.shapes.title.text = "Regional Performance Dashboard"
            slide1.placeholders[1].text = "Q3 2024 Summary"

            # Slide 2: Table slide
            slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
            slide2.shapes.title.text = "Revenue by Region"

            rows, cols = 6, 5
            left, top, width, height = Inches(0.5), Inches(1.5), Inches(9), Inches(4.5)
            table = slide2.shapes.add_table(rows, cols, left, top, width, height).table

            headers = ["Region", "Q3 Revenue ($M)", "vs Q2 ($M)", "YoY Growth (%)", "Target Hit?"]
            data = [
                ("North America", "142.3", "+8.2", "+16%", "Yes"),
                ("Europe", "87.6", "+3.1", "+9%", "Yes"),
                ("Asia Pacific", "54.2", "+6.8", "+22%", "Yes"),
                ("Latin America", "18.9", "-1.2", "+4%", "No"),
                ("Middle East & Africa", "12.1", "+0.5", "+11%", "Yes"),
            ]

            for c, header in enumerate(headers):
                cell = table.cell(0, c)
                cell.text = header
                cell.text_frame.paragraphs[0].runs[0].font.bold = True
                cell.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(255, 255, 255)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(31, 78, 121)

            for r, row_data in enumerate(data, start=1):
                for c, val in enumerate(row_data):
                    table.cell(r, c).text = val

            # Slide 3: Product table
            slide3 = prs.slides.add_slide(prs.slide_layouts[5])
            slide3.shapes.title.text = "Top 5 Products by Revenue"
            table2 = slide3.shapes.add_table(6, 4, Inches(0.5), Inches(1.5), Inches(9), Inches(4)).table
            h2 = ["Product", "Units Sold", "Revenue ($M)", "Margin (%)"]
            p2 = [
                ("Enterprise Suite", "1,240", "68.4", "45%"),
                ("Cloud Platform", "3,820", "57.3", "62%"),
                ("Analytics Module", "2,150", "38.7", "52%"),
                ("Mobile SDK", "8,900", "22.1", "71%"),
                ("Professional Services", "N/A", "18.6", "38%"),
            ]
            for c, h in enumerate(h2):
                cell = table2.cell(0, c)
                cell.text = h
                cell.text_frame.paragraphs[0].runs[0].font.bold = True
            for r, row_data in enumerate(p2, start=1):
                for c, val in enumerate(row_data):
                    table2.cell(r, c).text = val

            prs.save(out / filename)
        return self._entry(filename, "medium",
                           ["pptx", "tables", "regional-data", "product-data"],
                           False, "PPTX with two table slides: regional revenue breakdown and top product table.")

    def _complex_layout(self, out: Path, dry_run: bool) -> dict:
        filename = "syn_complex_layout.pptx"
        if dry_run:
            print(f"  [dry-run] would create {self.category}/{filename}")
        else:
            prs = Presentation()

            # Slide with multiple text boxes positioned manually
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

            # Add title text box
            txb = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9.4), Inches(0.8))
            tf = txb.text_frame
            tf.text = "Strategic Initiative Overview"
            tf.paragraphs[0].runs[0].font.size = Pt(28)
            tf.paragraphs[0].runs[0].font.bold = True
            tf.paragraphs[0].runs[0].font.color.rgb = RGBColor(31, 78, 121)

            # Left box
            left_box = slide.shapes.add_textbox(Inches(0.3), Inches(1.2), Inches(3), Inches(5))
            lf = left_box.text_frame
            lf.word_wrap = True
            lf.text = "Problem Statement"
            lf.paragraphs[0].runs[0].font.bold = True
            lf.paragraphs[0].runs[0].font.size = Pt(14)
            p = lf.add_paragraph()
            p.text = ("Current manual processes create bottlenecks in reporting. "
                      "Teams spend 40% of time on data aggregation rather than analysis.")
            p.runs[0].font.size = Pt(11)

            # Middle box
            mid_box = slide.shapes.add_textbox(Inches(3.5), Inches(1.2), Inches(3), Inches(5))
            mf = mid_box.text_frame
            mf.word_wrap = True
            mf.text = "Proposed Solution"
            mf.paragraphs[0].runs[0].font.bold = True
            mf.paragraphs[0].runs[0].font.size = Pt(14)
            p2 = mf.add_paragraph()
            p2.text = ("Implement unified data platform with automated pipelines. "
                       "Self-service dashboards eliminate manual report building.")
            p2.runs[0].font.size = Pt(11)

            # Right box
            right_box = slide.shapes.add_textbox(Inches(6.7), Inches(1.2), Inches(3), Inches(5))
            rf = right_box.text_frame
            rf.word_wrap = True
            rf.text = "Expected Outcomes"
            rf.paragraphs[0].runs[0].font.bold = True
            rf.paragraphs[0].runs[0].font.size = Pt(14)
            p3 = rf.add_paragraph()
            p3.text = ("75% reduction in report generation time. ROI: 340% over 3 years. "
                       "Team freed to focus on strategic analysis.")
            p3.runs[0].font.size = Pt(11)

            # Add a second slide with multiple shapes
            slide2 = prs.slides.add_slide(prs.slide_layouts[6])
            # Timeline-style text boxes
            for i, (label, detail) in enumerate([
                ("Phase 1\nMar–Jun", "Infrastructure Setup\nData lake, ETL pipelines, security review"),
                ("Phase 2\nJul–Sep", "Core Analytics\nDashboard framework, first 5 reports live"),
                ("Phase 3\nOct–Dec", "Advanced Features\nML models, anomaly detection, self-service"),
            ]):
                box = slide2.shapes.add_textbox(
                    Inches(0.3 + i * 3.3), Inches(1.5), Inches(3), Inches(3.5)
                )
                tf_ph = box.text_frame
                tf_ph.word_wrap = True
                tf_ph.text = label
                tf_ph.paragraphs[0].runs[0].font.bold = True
                tf_ph.paragraphs[0].runs[0].font.size = Pt(16)
                tf_ph.paragraphs[0].runs[0].font.color.rgb = RGBColor(31, 78, 121)
                pd = tf_ph.add_paragraph()
                pd.text = detail
                pd.runs[0].font.size = Pt(11)

            prs.save(out / filename)
        return self._entry(filename, "medium",
                           ["pptx", "multiple-text-boxes", "complex-layout", "strategic-overview"],
                           False, "PPTX with manually positioned multiple text boxes per slide — tests non-standard layout extraction.")

    def _speaker_notes(self, out: Path, dry_run: bool) -> dict:
        filename = "syn_speaker_notes.pptx"
        if dry_run:
            print(f"  [dry-run] would create {self.category}/{filename}")
        else:
            prs = Presentation()

            slides_content = [
                (
                    "Introduction to Machine Learning",
                    "Foundational concepts for business leaders",
                    ("Machine learning (ML) is a branch of AI that enables systems to learn "
                     "from data without being explicitly programmed. Key types include supervised "
                     "learning (labeled data), unsupervised learning (pattern discovery), and "
                     "reinforcement learning (reward-based optimization). Today we'll focus on "
                     "practical business applications and implementation considerations."),
                ),
                (
                    "Why ML Matters for Your Business",
                    ["Automate repetitive decision-making at scale",
                     "Uncover patterns invisible to human analysts",
                     "Personalize customer experiences in real time",
                     "Predict failures before they occur",
                     "Optimize pricing and inventory dynamically"],
                    ("Emphasize ROI examples: Netflix saves $1B/year via recommendations. "
                     "Amazon's demand forecasting reduces stockouts by 30%. "
                     "GE's predictive maintenance cuts downtime by 20%. "
                     "Pause here for questions — audience often skeptical at this slide."),
                ),
                (
                    "Common Use Cases by Industry",
                    ["Healthcare: diagnosis assistance, drug discovery, patient readmission prediction",
                     "Finance: fraud detection, credit scoring, algorithmic trading",
                     "Retail: demand forecasting, customer segmentation, price optimization",
                     "Manufacturing: quality control, predictive maintenance, supply chain optimization"],
                    ("Walk through one healthcare and one finance example in detail. "
                     "The healthcare examples resonate well with this audience. "
                     "Do NOT mention autonomous vehicles — off-topic and derails discussion."),
                ),
                (
                    "Getting Started: Practical Roadmap",
                    ["Step 1: Identify a high-value, well-defined problem (3 months)",
                     "Step 2: Audit your data quality and availability (2 months)",
                     "Step 3: Run a focused proof of concept (3 months)",
                     "Step 4: Evaluate results and build business case (1 month)",
                     "Step 5: Scale and operationalize (6-12 months)"],
                    ("This is the most important slide. Leave 10 minutes for discussion here. "
                     "Key message: start small, prove value, then scale. "
                     "Typical failure mode: trying to do everything at once. "
                     "Have the maturity assessment framework ready as a follow-up resource."),
                ),
            ]

            for title, content, notes_text in slides_content:
                layout = prs.slide_layouts[1] if isinstance(content, list) else prs.slide_layouts[1]
                slide = prs.slides.add_slide(layout)
                slide.shapes.title.text = title

                if isinstance(content, list):
                    tf = slide.placeholders[1].text_frame
                    tf.text = content[0]
                    for item in content[1:]:
                        p = tf.add_paragraph()
                        p.text = item
                else:
                    slide.placeholders[1].text = content

                # Add speaker notes
                notes_slide = slide.notes_slide
                notes_tf = notes_slide.notes_text_frame
                notes_tf.text = notes_text

            prs.save(out / filename)
        return self._entry(filename, "medium",
                           ["pptx", "speaker-notes", "ml-overview", "extensive-notes"],
                           False, "Four-slide PPTX with extensive speaker notes on each slide — tests notes extraction.")
